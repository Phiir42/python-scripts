from __future__ import annotations

import math
import os
from typing import Any, Dict, List, MutableMapping, Optional, Tuple

import numpy as np
import numpy.typing as npt
from matplotlib.backends.backend_pdf import PdfPages
from scipy.interpolate import make_interp_spline
from scipy.signal import find_peaks

try:
    # Types for lmfit objects are not available; use Any in annotations.
    from lmfit import Minimizer, Parameters
except ImportError as e:
    raise ImportError("Please `pip install lmfit` to enable peak fitting.") from e


# ----------------------------
# Debug capture of plots (PNG + multipage PDF + optional PPTX)
# ----------------------------
_DEBUG_CAPTURE: Dict[str, Any] = {
    "save_dir": None,  # type: Optional[str]
    "pdf": None,  # type: Optional[PdfPages]
    "paths": [],  # type: List[str]
}


def start_debug_capture(save_dir: str, pdf_name: str = "fits.pdf") -> None:
    """Begin capturing debug plot outputs."""
    os.makedirs(save_dir, exist_ok=True)
    pdf_path = os.path.join(save_dir, pdf_name)
    _DEBUG_CAPTURE["save_dir"] = save_dir
    _DEBUG_CAPTURE["pdf"] = PdfPages(pdf_path)
    _DEBUG_CAPTURE["paths"] = []


def finish_debug_capture(make_pptx: bool = True, pptx_name: str = "fits.pptx") -> None:
    """Finish capturing: close PDF and optionally build a PPTX from saved PNGs."""
    # close PDF
    if _DEBUG_CAPTURE["pdf"] is not None:
        try:
            _DEBUG_CAPTURE["pdf"].close()
        except Exception:
            pass
        _DEBUG_CAPTURE["pdf"] = None

    # build PPTX if requested
    if make_pptx and _DEBUG_CAPTURE["paths"]:
        try:
            from pptx import Presentation  # type: ignore[import-not-found]
            from pptx.util import Inches  # type: ignore[import-not-found]

            prs = Presentation()
            blank = prs.slide_layouts[6]
            for p in _DEBUG_CAPTURE["paths"]:
                slide = prs.slides.add_slide(blank)
                slide.shapes.add_picture(p, Inches(0.3), Inches(0.3), width=Inches(9.0))
            out_pptx = os.path.join(_DEBUG_CAPTURE["save_dir"], pptx_name)
            prs.save(out_pptx)
        except ImportError:
            # PPTX generation is optional; ignore if python-pptx isn't installed.
            pass

    _DEBUG_CAPTURE["save_dir"] = None
    _DEBUG_CAPTURE["paths"] = []


# ----------------------------
# Chi-square accumulator (for batch debug summaries)
# ----------------------------
_CHI2_SUM: float = 0.0
_CHI2_COUNT: int = 0
_CHI2_NONCONV: int = 0
_CHI2_ITEMS: List[Dict[str, Any]] = []  # optional detail records


def chi2_reset() -> None:
    """Reset the running χ² counters (sum, count, non-converged)."""
    global _CHI2_SUM, _CHI2_COUNT, _CHI2_NONCONV, _CHI2_ITEMS
    _CHI2_SUM = 0.0
    _CHI2_COUNT = 0
    _CHI2_NONCONV = 0
    _CHI2_ITEMS = []


def chi2_add(
    series_label: str,
    droplet_id: int,
    redchi: Any,
    success: bool,
    strategy: str = "",
) -> None:
    """
    Add one fit's χ² to the running totals.

    Parameters
    ----------
    series_label : str
        Folder / sample label (e.g., 'AD3d_...').
    droplet_id : int
        Lipid ID for tracing.
    redchi : float-like
        Reduced chi-square reported by the fit.
    success : bool
        Optimizer success flag (True/False).
    strategy : str
        Which strategy was used, e.g. 'data-seeded + x8'.
    """
    global _CHI2_SUM, _CHI2_COUNT, _CHI2_NONCONV, _CHI2_ITEMS
    try:
        rc = float(redchi)  # type: ignore[arg-type]
    except Exception:
        rc = float("nan")

    if not math.isfinite(rc):
        rc = float("nan")

    if math.isfinite(rc):
        _CHI2_SUM += rc
        _CHI2_COUNT += 1
    else:
        # still count the attempt, but record as NaN in the details
        _CHI2_COUNT += 1

    if not success:
        _CHI2_NONCONV += 1

    _CHI2_ITEMS.append(
        {
            "series": series_label,
            "lipid_id": droplet_id,
            "redchi": rc,
            "success": bool(success),
            "strategy": strategy or "",
        }
    )


def report_chi2_summary() -> Dict[str, Any]:
    """
    Print a one-line and a detailed χ² summary. Returns a dict with the stats.
    """
    n = _CHI2_COUNT if _CHI2_COUNT > 0 else 1
    mean = _CHI2_SUM / (_CHI2_COUNT if _CHI2_COUNT else 1)
    print(
        f"[χ²] total={_CHI2_SUM:.6g}  mean={mean:.6g}  "
        f"n={_CHI2_COUNT}  non-converged={_CHI2_NONCONV}"
    )
    # Optional: top 5 worst by χ² (finite only)
    finite = [d for d in _CHI2_ITEMS if math.isfinite(d["redchi"])]
    worst = sorted(finite, key=lambda d: float(d["redchi"]), reverse=True)[:5]
    if worst:
        print("[χ²] worst 5:")
        for d in worst:
            print(
                "      series="
                f"{str(d['series']):20s}  lipid={int(d['lipid_id']):>4}  "
                f"χ²={float(d['redchi']):.6g}  success={bool(d['success'])}  "
                f"{str(d['strategy'])}"
            )
    return {
        "sum": _CHI2_SUM,
        "mean": mean,
        "n": n,
        "non_converged": _CHI2_NONCONV,
        "items": list(_CHI2_ITEMS),
    }


# ----------------------------
# Core model
# ----------------------------
def _cars_model_complex(
    x: npt.ArrayLike, par: MutableMapping[str, Any]
) -> np.ndarray:
    """
    Complex-sum CARS model with optional broad 3010 band and a positive
    strictly linear supporting baseline.
    Returns: |sum(...)|^2 + baseline(x)
    """
    x_arr = np.asarray(x, dtype=float)

    def _p(name: str, default: float = 0.0) -> float:
        v = par[name] if (hasattr(par, "__getitem__") and name in par) else default
        return float(getattr(v, "value", v))

    cc = np.zeros_like(x_arr, dtype=np.complex128)

    # 7 narrow bands + optional broad (x8 is mandatory in defaults)
    for k in range(1, 8 + int(_p("enable_x8", 0))):
        A = _p(f"A{k}")
        x0 = _p(f"x{k}")
        w = _p(f"w{k}")
        cc += A / ((x_arr - x0) + 1j * w)

    # Complex non-resonant background (constant + optional linear slope, both complex)
    cc += _p("ANR_re", 0.0) + 1j * _p("ANR_im", 0.0)
    xref = _p("xref", 3000.0)
    if _p("enable_linear_nr", 0):
        cc += (_p("BNR1_re", 0.0) + 1j * _p("BNR1_im", 0.0)) * (x_arr - xref)

    # Far off-resonant pseudo-terms (complex Lorentzians far from CH band)
    for tag in ("OR1", "OR2"):
        A = _p(f"A{tag}", 0.0)
        x0 = _p(f"x{tag}", 0.0)
        w = _p(f"w{tag}", 20.0)
        cc += A / ((x_arr - x0) + 1j * w)

    y = np.abs(cc) ** 2

    # --- strictly linear baseline ---
    x_arr_min = float(np.min(x_arr))
    # NEW: anchor the line at (x_arr_min, BL_left) and (BL_xright, BL_right)
    BL_left   = _p("BL_left",  0.0)
    BL_right  = _p("BL_right", 0.0)
    BL_xright = _p("BL_xright", float(np.max(x_arr)))  # default to end if not provided
    den = max(BL_xright - x_arr_min, 1e-9)
    slope = (BL_right - BL_left) / den
    baseline = BL_left + slope * (x_arr - x_arr_min)
    y = y + baseline
    return y


def _residual(
    par: MutableMapping[str, Any],
    x: npt.ArrayLike,
    y: npt.ArrayLike,
    w: Optional[npt.ArrayLike] = None,
) -> np.ndarray:
    r = _cars_model_complex(x, par) - np.asarray(y, dtype=float)
    if w is None:
        return r
    return r * np.asarray(w, dtype=float)


# ----------------------------
# Parameter seeding
# ----------------------------
def _make_default_params_from_config(config: Dict[str, Any]) -> Parameters:
    """
    Seeds/bounds. If config has no 'peak_fit' key, robust defaults are used.
    """
    pf: Dict[str, Any] = config.get("peak_fit", {})

    centers_seed = pf.get(
        "centers_seed", [2859.9, 2884.7, 2936.0, 3020.0, 2812.0, 2908.6, 2967.7]
    )
    centers_low = pf.get("centers_low", [2845, 2875, 2925, 3012, 2805, 2900, 2950])
    centers_high = pf.get("centers_high", [2865, 2895, 2943, 3026, 2825, 2920, 3000])

    amps_seed = pf.get("amps_seed", [13.92, 8.84, 7.08, 0.50, 1.25, 9.77, -1.41])
    amp_min = pf.get("amp_min", None)
    amp_max = pf.get("amp_max", None)

    # Wider mins for peaks that tend to needle (1,2,6). Allow w4 to be broader.
    widths_seed = pf.get("widths_seed", [10, 12, 14.8, 14, 20, 12, 22])
    widths_low = pf.get("widths_low", [8, 8, 6, 8, 5, 8, 6])
    widths_high = pf.get("widths_high", [30, 30, 30, 60, 30, 40, 30])

    p = Parameters()
    for i in range(7):
        p.add(
            f"x{i+1}",
            value=centers_seed[i],
            min=centers_low[i],
            max=centers_high[i],
            vary=True,
        )
        p.add(f"A{i+1}", value=amps_seed[i], min=amp_min, max=amp_max, vary=True)
        p.add(
            f"w{i+1}",
            value=widths_seed[i],
            min=widths_low[i],
            max=widths_high[i],
            vary=True,
        )
    # ---- Robust pruning: freeze selected peaks without min==max ----
    # e.g. config["peak_fit"]["drop_peaks"] = [5]   # to drop Peak 5
    drop_set = set(int(k) for k in pf.get("drop_peaks", []))
    for i in range(7):
        idx = i + 1
        if idx in drop_set:
            # Freeze x and w
            if f"x{idx}" in p:
                p[f"x{idx}"].set(vary=False)
            if f"w{idx}" in p:
                p[f"w{idx}"].set(vary=False)

            # Freeze amplitude to a feasible value WITHOUT changing min/max
            if f"A{idx}" in p:
                A_par = p[f"A{idx}"]
                # Prefer 0 if it's within bounds, else pick a bound safely inside
                A_fix = 0.0
                if A_par.min is not None and A_fix < A_par.min:
                    A_fix = float(A_par.min) + 1e-12
                if A_par.max is not None and A_fix > A_par.max:
                    A_fix = float(A_par.max) - 1e-12
                A_par.set(value=A_fix, vary=False)

    # Complex NR background
    p.add(
        "ANR_re",
        value=pf.get("ANR_re_seed", 0.02),
        min=pf.get("ANR_re_min", -2.0),
        max=pf.get("ANR_re_max", 2.0),
        vary=pf.get("ANR_re_vary", True),
    )
    p.add(
        "ANR_im",
        value=pf.get("ANR_im_seed", 0.00),
        min=pf.get("ANR_im_min", -2.0),
        max=pf.get("ANR_im_max", 2.0),
        vary=pf.get("ANR_im_vary", True),
    )

    # Optional linear complex NR slope (ENABLED by default so wings can tilt)
    p.add("enable_linear_nr", value=pf.get("enable_linear_nr_seed", 1), vary=False)
    p.add("xref", value=pf.get("xref", 3000.0), vary=False)
    p.add(
        "BNR1_re",
        value=pf.get("BNR1_re_seed", 0.0),
        min=-2e-2,
        max=2e-2,
        vary=pf.get("BNR1_re_vary", True),
    )
    p.add(
        "BNR1_im",
        value=pf.get("BNR1_im_seed", 0.0),
        min=-2e-2,
        max=2e-2,
        vary=pf.get("BNR1_im_vary", True),
    )

    # Far off-resonant pseudo-term: let it actually move
    p.add(
        "xOR1",
        value=pf.get("xOR1_seed", 3549.0),
        min=pf.get("xOR1_min", 3200),
        max=pf.get("xOR1_max", 4200),
        vary=True,
    )
    p.add(
        "AOR1",
        value=pf.get("AOR1_seed", 0.0),
        min=pf.get("AOR1_min", -2.0),
        max=pf.get("AOR1_max", 2.0),
        vary=True,
    )
    p.add(
        "wOR1",
        value=pf.get("wOR1_seed", 40.0),
        min=pf.get("wOR1_min", 10),
        max=pf.get("wOR1_max", 100),
        vary=True,
    )

    # Second pseudo-term still off by default
    p.add("xOR2", value=pf.get("xOR2_seed", 0.0), min=0, max=1500, vary=False)
    p.add("AOR2", value=pf.get("AOR2_seed", 0.0), vary=False)
    p.add("wOR2", value=pf.get("wOR2_seed", 20.0), min=5, max=60, vary=False)

    # Make the broad 3010 component mandatory and distinct from x4 (narrower)
    p.add("enable_x8", value=1, vary=False)
    p.add("x8", value=pf.get("x8_seed", 3010.0), min=3000.0, max=3028.0, vary=True)
    p.add("A8", value=pf.get("A8_seed", 0.20), min=-3.0, max=3.0, vary=True)
    p.add("w8", value=pf.get("w8_seed", 50.0), min=30.0, max=100.0, vary=True)

    # NEW: strictly positive *linear* baseline at [left, right]
    # Midpoint removed; 0.3 default ensures modest uplift at edges.
    p.add("BL_left",   value=pf.get("BL_left_seed",   0.3), min=-1.0, max=3.0, vary=False)
    p.add("BL_right",  value=pf.get("BL_right_seed",  0.3), min=-1.0, max=3.0, vary=False)
    # NEW: x-position (cm^-1) of the right baseline anchor; fixed (no fitting).
    p.add("BL_xright", value=pf.get("BL_xright_seed", 3000.0), vary=False)

    return p


def _seed_centers_from_data(
    x: npt.ArrayLike,
    y: npt.ArrayLike,
    target_centers: Tuple[int, ...] = (2855, 2890, 2935, 3019, 2812, 2909, 2968),
) -> List[float]:
    # local smoothing for robust peak finding
    from skimage.filters import gaussian  # local import; skimage is optional at runtime

    x_arr = np.asarray(x, dtype=float)
    y_arr = np.asarray(y, dtype=float)
    y_s = gaussian(y_arr, sigma=1, preserve_range=True)
    prom = float(np.percentile(y_s, 60) / 6.0)
    pk_idx, _ = find_peaks(y_s, prominence=prom, distance=2)
    pk_x = x_arr[pk_idx]
    centers: List[float] = []
    for t in target_centers:
        val = float(t if len(pk_x) == 0 else pk_x[np.argmin(np.abs(pk_x - t))])
        centers.append(val)
    return centers


# ----------------------------
# Baseline: supporting-line under the data
# ----------------------------
def _supporting_line_under_curve(
    x: npt.ArrayLike,
    y_raw: npt.ArrayLike,
    tol_frac: float = 0.005,
    smooth_sigma: float = 0.0,
) -> Tuple[float, float]:
    """
    Return slope m and intercept c of the *highest* line y = m*x + c
    that never exceeds the (optionally smoothed) raw spectrum:
        m*x_i + c <= y_i  for all i
    'tol_frac' gives a tiny slack (fraction of dynamic range) to avoid
    rejecting nearly-tangent candidates due to noise.
    """
    x_arr = np.asarray(x, dtype=float)
    y_arr = np.asarray(y_raw, dtype=float)
    finite = np.isfinite(y_arr)
    x_f = x_arr[finite]
    y_f = y_arr[finite]
    if x_f.size < 2:
        return 0.0, float(np.nanmin(y_f) if y_f.size else 0.0)

    # optional gentle smoothing to mitigate single-point noise
    if smooth_sigma > 0:
        try:
            from skimage.filters import gaussian  # optional
            y_check = gaussian(y_f, sigma=smooth_sigma, preserve_range=True)
        except Exception:
            y_check = y_f
    else:
        y_check = y_f

    ymin = float(np.nanmin(y_f))
    ymax = float(np.nanmax(y_f))
    dyn  = float(max(ymax - ymin, 1.0))
    tol  = tol_frac * dyn

    # Brute-force over all pairs; keep the feasible line with the largest sum
    best_score = -np.inf
    best_m, best_c = 0.0, ymin

    n = x_f.size
    for i in range(n - 1):
        xi, yi = x_f[i], y_f[i]
        for j in range(i + 1, n):
            xj, yj = x_f[j], y_f[j]
            if xj == xi:
                continue
            m = (yj - yi) / (xj - xi)
            c = yi - m * xi

            # Feasibility: line never above spectrum (allow tiny slack tol)
            if np.any(m * x_f + c > (y_check + tol)):
                continue

            # Score: choose the "highest" feasible line (largest total baseline)
            score = float(np.sum(m * x_f + c))
            if score > best_score:
                best_score = score
                best_m, best_c = m, c

    return best_m, best_c


# ----------------------------
# Fitting utilities
# ----------------------------
def _weights_for_fit(x: npt.ArrayLike, pf: Dict[str, Any]) -> np.ndarray:
    # gentle weighting: boost endpoints and the 2880/2930 shoulders
    x_arr = np.asarray(x, dtype=float)
    w = np.ones_like(x_arr, dtype=float)

    def _bump(lo: float, hi: float, boost: float) -> None:
        if boost <= 1.0:
            return
        c = 0.5 * (lo + hi)
        r = 0.5 * (hi - lo)
        band = np.clip((r - np.abs(x_arr - c)) / r, 0.0, 1.0)
        w[:] *= 1.0 + (boost - 1.0) * band

    boosts = pf if isinstance(pf, dict) else {}
    _bump(2780.0, 2820.0, float(boosts.get("edge_lo_boost", 1.3)))
    _bump(2990.0, 3030.0, float(boosts.get("edge_hi_boost", 1.3)))
    _bump(2870.0, 2895.0, float(boosts.get("mid_2880_boost", 1.15)))
    _bump(2925.0, 2945.0, float(boosts.get("mid_2930_boost", 1.15)))
    return w


def _run_once(
    params: Parameters,
    x: npt.ArrayLike,
    y: npt.ArrayLike,
    w: Optional[npt.ArrayLike] = None,
    loss: str = "soft_l1",
) -> Any:
    minner = Minimizer(_residual, params, fcn_args=(x, y, w))
    # Disable covariance calculation to avoid RuntimeWarnings when the Hessian
    # is ill-conditioned or when using robust loss (soft_l1/huber).
    return minner.least_squares(
        method="trf",
        max_nfev=10000,
        loss=loss,
        f_scale=0.1,
    )


def _jitter_params(p: Parameters, frac: float = 0.05) -> Parameters:
    """Return a shallow-copied lmfit.Parameters with values jittered within bounds."""
    q = p.copy()
    rng = np.random.default_rng()
    for _name, par in q.items():
        if not par.vary:
            continue
        v = float(par.value)
        lo = float(par.min) if par.min is not None else v - abs(v) - 1.0
        hi = float(par.max) if par.max is not None else v + abs(v) + 1.0
        span = hi - lo
        if not np.isfinite(span) or span <= 0:
            continue
        dv = frac * span * float(rng.standard_normal())
        newv = np.clip(v + dv, lo, hi)
        par.set(value=float(newv))
    return q


def _fit_with_retries(
    base_params: Parameters,
    x: npt.ArrayLike,
    y: npt.ArrayLike,
    w: Optional[npt.ArrayLike],
    retries: int = 4,
) -> Tuple[Any, Parameters]:
    """Try base + several jittered restarts; return best result & params."""
    best: Optional[Any] = None
    best_params: Parameters = base_params

    # try two loss types for stubborn cases
    for loss in ("soft_l1", "huber"):
        # base try
        res = _run_once(base_params, x, y, w, loss=loss)
        cand = (
            float(res.redchi)
            if np.isfinite(getattr(res, "redchi", np.inf))
            else float("inf")
        )
        best = res
        best_params = base_params

        # jittered restarts
        for _ in range(retries):
            pj = _jitter_params(base_params, frac=0.08)
            rj = _run_once(pj, x, y, w, loss=loss)
            rc = (
                float(rj.redchi)
                if np.isfinite(getattr(rj, "redchi", np.inf))
                else float("inf")
            )
            if rc < cand:
                cand, best, best_params = rc, rj, pj

    # --- POLISH STEP: one final pass with 'linear' loss from the best params ---
    # This tightens the solution and often returns a "converged" termination.
    from lmfit import Minimizer  # local import to keep top clean
    pol_min = Minimizer(_residual, best.params, fcn_args=(x, y, w))
    pol = pol_min.least_squares(
        method="trf",
        loss="linear",     # pure least squares (no robust weighting)
        f_scale=1.0,
        max_nfev=1500,     # short polish; adjust if you like
    )
    pol_redchi = float(getattr(pol, "redchi", np.inf))
    best_redchi = float(getattr(best, "redchi", np.inf))
    if np.isfinite(pol_redchi) and (pol_redchi <= best_redchi):
        best = pol
        best_params = pol.params

    # best is always set (at least base try)
    assert best is not None
    return best, best_params


# ----------------------------
# Public API
# ----------------------------
def fit_cars_peaks(
    wavenumbers_cm1: npt.ArrayLike, intensity: npt.ArrayLike, config: Dict[str, Any]
) -> Dict[str, Any]:
    """
    Fit CH-band CARS spectrum with:
      - 7 narrow resonances + mandatory broad 3010 component (x8)
      - complex NR (constant + optional linear slope)
      - far off-resonant pseudo term(s)
      - positive piecewise-linear baseline (BL_left / BL_right at edges/xref)
      - multi-start with jitter for robust convergence
    """
    x = np.asarray(wavenumbers_cm1, dtype=float)
    y_raw = np.asarray(intensity, dtype=float)
    finite = np.isfinite(y_raw)
    if not np.any(finite):
        raise ValueError("No finite intensity values in spectrum.")

    # per-spectrum offset and dynamic range from FINITE samples only
    ymin = float(np.nanmin(y_raw[finite]))
    ymax = float(np.nanmax(y_raw[finite]) - ymin)
    if not np.isfinite(ymax) or ymax <= 0:
        ymax = 1.0

    # --- Seed baseline as the *supporting line* under the raw spectrum ---
    pf = config.get("peak_fit", {})
    m_sup, c_sup = _supporting_line_under_curve(
        x, y_raw,
        tol_frac=float(pf.get("baseline_tol_frac", 0.005)),
        smooth_sigma=float(pf.get("baseline_smooth_sigma", 0.0)),
    )
    
    # Values at the measurement endpoints
    x_min = float(np.min(x))
    x_max = float(np.max(x))
    bl_left_raw  = m_sup * x_min + c_sup
    bl_right_raw = m_sup * x_max + c_sup
    
    # Normalize into the [0,1]-like working scale the fitter uses
    bl_left_seed  = (bl_left_raw  - ymin) / ymax
    bl_right_seed = (bl_right_raw - ymin) / ymax

    # normalized working array for fitting (leave NaNs out of the way)
    y = (y_raw - ymin) / ymax
    # replace remaining non-finite with local linear interpolation; if that fails, set to 0
    if np.any(~finite):
        try:
            y[~finite] = np.interp(x[~finite], x[finite], y[finite])
        except Exception:
            y[~finite] = 0.0
            
    weights = _weights_for_fit(x, pf.get("weights", {}))

    # Use a copy of config so we don't mutate the caller's dict
    cfg_for_fit = dict(config)
    cfg_for_fit_pf = dict(cfg_for_fit.get("peak_fit", {}))
    cfg_for_fit_pf.setdefault("drop_peaks", [5])
    cfg_for_fit_pf["BL_left_seed"]   = float(bl_left_seed)
    cfg_for_fit_pf["BL_right_seed"]  = float(bl_right_seed)
    cfg_for_fit_pf["BL_xright_seed"] = float(x_max)   # line spans the full domain
    cfg_for_fit["peak_fit"] = cfg_for_fit_pf

    def _pack_output(res: Any, params: Parameters, label: str) -> Dict[str, Any]:
        out: Dict[str, Any] = {f"x{k}": params[f"x{k}"].value for k in range(1, 8)}
        out.update({f"A{k}": params[f"A{k}"].value for k in range(1, 8)})
        out.update({f"w{k}": params[f"w{k}"].value for k in range(1, 8)})

        # extras
        for nm in (
            "ANR_re",
            "ANR_im",
            "BNR1_re",
            "BNR1_im",
            "xref",
            "xOR1",
            "AOR1",
            "wOR1",
            "xOR2",
            "AOR2",
            "wOR2",
            "x8",
            "A8",
            "w8",
            "BL_left",
            "BL_right",
            "BL_xright",
        ):
            if nm in params:
                out[nm] = params[nm].value

        # diag
        msg = str(getattr(res, "message", "")).lower()
        redchi_val = float(getattr(res, "redchi", float("nan")))
        raw_success = bool(getattr(res, "success", False))
        heuristic_ok = (
            (math.isfinite(redchi_val) and redchi_val < 2e-2)  # very good χ²
            or ("gtol" in msg or "xtol" in msg or "ftol" in msg)  # typical terminations
        )
    
        out.update(
            {
                "success": raw_success or heuristic_ok,
                "redchi": redchi_val,
                "nfev": float(getattr(res, "nfev", float("nan"))),
                "message": str(getattr(res, "message", "")),
                "active_bounds": list(getattr(res, "active_mask", [])),
                "strategy_used": label,
            }
        )

        try:
            out["y_fit"] = _cars_model_complex(x, params)
        except Exception:
            out["y_fit"] = None
        out["_scale"] = float(ymax)
        out["_offset"] = float(ymin)
        return out

    # Strategy A: defaults + restarts
    pA = _make_default_params_from_config(cfg_for_fit)
    resA, _ = _fit_with_retries(pA, x, y, weights, retries=4)
    outA = _pack_output(resA, resA.params, "default-seeds+restarts")

    # Early accept if excellent
    if outA["success"] and np.isfinite(outA["redchi"]) and outA["redchi"] < 1e-2:
        return outA

    # Strategy B: data-seeded centers + restarts
    pB = _make_default_params_from_config(cfg_for_fit)
    new_centers = _seed_centers_from_data(x, y)
    for i, c in enumerate(new_centers, 1):
        lo = pB[f"x{i}"].min if pB[f"x{i}"].min is not None else (c - 20)
        hi = pB[f"x{i}"].max if pB[f"x{i}"].max is not None else (c + 20)
        pB[f"x{i}"].set(value=float(np.clip(c, lo, hi)))
    resB, _ = _fit_with_retries(pB, x, y, weights, retries=4)
    outB = _pack_output(resB, resB.params, "data-seeded+restarts")

    # Choose better of A/B
    if outB["success"] and (
        not np.isfinite(outA["redchi"]) or outB["redchi"] <= outA["redchi"]
    ):
        return outB
    return outA


# ----------------------------
# Debug plot
# ----------------------------
def _plot_peak_fit_debug(
    x_cm1: npt.ArrayLike,
    y_vals: npt.ArrayLike,
    fit_result: Dict[str, Any],
    droplet_id: Any,
    category: str,
    location: str,
    marker: str,
) -> None:
    """
    Plot raw hyperspectrum, individual components, and the total fit.
    The red 'Total fit' includes resonances, complex NR (and its linear slope
    if enabled), far off-resonant pseudo terms, and the additive baseline BL0/BL1.
    """
    import matplotlib.pyplot as plt  # local to avoid import if not plotting

    x_arr = np.asarray(x_cm1, dtype=float)
    y_arr = np.asarray(y_vals, dtype=float)
    scale = float(fit_result.get("_scale", 1.0))
    offset = float(fit_result.get("_offset", 0.0))

    x_dense = np.linspace(float(np.min(x_arr)), float(np.max(x_arr)), 400)
    try:
        y_smooth = make_interp_spline(x_arr, y_arr, k=3)(x_dense)
    except Exception:
        y_smooth = np.interp(x_dense, x_arr, y_arr)

    plt.figure(figsize=(7, 5))
    plt.plot(x_dense, y_smooth, "-", color="black", lw=1.5, label="Raw (spline)")
    plt.scatter(x_arr, y_arr, c="k", s=25, marker="o", label="Raw points")

    # Total fit (prefer precomputed y_fit which already includes baseline terms)
    y_fit_norm = fit_result.get("y_fit", None)
    if y_fit_norm is not None:
        y_fit_norm_arr = np.asarray(y_fit_norm, dtype=float)
        try:
            y_fit_dense = make_interp_spline(x_arr, y_fit_norm_arr, k=3)(x_dense)
        except Exception:
            y_fit_dense = np.interp(x_dense, x_arr, y_fit_norm_arr)
        # convert back to raw units
        plt.plot(x_dense, y_fit_dense * scale + offset, "r-", lw=2, label="Total fit")

        # linear baseline only
        BL_left   = float(fit_result.get("BL_left", 0.0))
        BL_right  = float(fit_result.get("BL_right", 0.0))
        BL_xright = float(fit_result.get("BL_xright", float(np.max(x_arr))))
        x_min = float(np.min(x_arr))
        den = max(BL_xright - x_min, 1e-9)
        slope = (BL_right - BL_left) / den
        baseline = (BL_left + slope * (x_dense - x_min)) * scale + offset
        if np.any(baseline != 0):
            plt.plot(
                x_dense,
                baseline,
                linestyle=(0, (3, 3)),
                lw=1.2,
                color="tab:red",
                label="Linear baseline",
            )

    # Individual peak components
    colors = plt.cm.tab10.colors
    ks = sorted(
        {
            int(k[1:])
            for k in fit_result.keys()
            if k.startswith("x") and k[1:].isdigit()
        }
    )
    for idx, k in enumerate(ks):
        Ak = float(fit_result.get(f"A{k}", 0.0))
        x0 = float(fit_result.get(f"x{k}", 0.0))
        wk = float(fit_result.get(f"w{k}", 0.0))
        if Ak == 0.0 or wk == 0.0:
            continue
        comp = np.abs(Ak / ((x_dense - x0) + 1j * wk)) ** 2
        plt.plot(
            x_dense,
            comp * scale,
            "--",
            lw=1.2,
            color=colors[(idx) % len(colors)],
            label=f"Peak {k} ({x0:.0f} cm⁻¹)",
        )

    ttl = f"LipidID {droplet_id} | {category}, {location}, {marker}"
    if "strategy_used" in fit_result:
        try:
            ttl += (
                f"  [{fit_result['strategy_used']}, "
                f"χ²≈{float(fit_result.get('redchi', np.nan)):.3g}]"
            )
        except Exception:
            ttl += f"  [{fit_result['strategy_used']}]"
    if fit_result.get("success", True) is False:
        ttl += "  (non-converged)"
    plt.xlabel("Raman shift (cm$^{-1}$)")
    plt.ylabel("Intensity (a.u.)")
    plt.title(ttl)
    plt.legend(fontsize=8, ncol=2)
    plt.tight_layout()

    # --- SAVE into capture if enabled (save BEFORE show/close) ---
    fig = plt.gcf()
    try:
        if _DEBUG_CAPTURE["pdf"] is not None:
            _DEBUG_CAPTURE["pdf"].savefig(fig, bbox_inches="tight")
        if _DEBUG_CAPTURE["save_dir"] is not None:
            # use lipid id in filename when available
            stem = (
                f"fit_{int(droplet_id):04d}"
                if str(droplet_id).isdigit()
                else f"fit_{droplet_id}"
            )
            png_path = os.path.join(_DEBUG_CAPTURE["save_dir"], f"{stem}.png")
            fig.savefig(png_path, dpi=200, bbox_inches="tight")
            _DEBUG_CAPTURE["paths"].append(png_path)
    except Exception:
        pass

    # Now optionally display, then close.
    plt.show()
    plt.close(fig)
